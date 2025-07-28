import os
import re
from typing import List, Dict, Any, Optional
from pathlib import Path
import logging

# Document processing libraries
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import pandas as pd

# LangChain components
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document as LangchainDocument

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """
    A comprehensive document processor that handles multiple file formats
    and provides intelligent text chunking for RAG applications.
    """

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the document processor.

        Args:
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""],
        )

        # Supported file extensions
        self.supported_extensions = {
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".doc": self._process_docx,
            ".pptx": self._process_pptx,
            ".txt": self._process_txt,
            ".png": self._process_image,
            ".jpg": self._process_image,
            ".jpeg": self._process_image,
            ".tiff": self._process_image,
            ".bmp": self._process_image,
            ".csv": self._process_csv,
        }

    def process_file(self, file_path: str) -> List[LangchainDocument]:
        """
        Process a single file and return chunked documents.

        Args:
            file_path: Path to the file to process

        Returns:
            List of LangchainDocument objects
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_extension = file_path.suffix.lower()

        if file_extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_extension}")

        try:
            logger.info(f"Processing file: {file_path}")
            text_content = self.supported_extensions[file_extension](file_path)

            if not text_content.strip():
                logger.warning(f"No text content extracted from {file_path}")
                return []

            # Create metadata
            metadata = {
                "source": str(file_path),
                "file_type": file_extension,
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
            }

            # Split text into chunks
            chunks = self.text_splitter.split_text(text_content)

            # Create LangchainDocument objects
            documents = []
            for i, chunk in enumerate(chunks):
                doc = LangchainDocument(
                    page_content=chunk,
                    metadata={**metadata, "chunk_id": i, "total_chunks": len(chunks)},
                )
                documents.append(doc)

            logger.info(f"Created {len(documents)} chunks from {file_path}")
            return documents

        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            raise

    def process_directory(self, directory_path: str) -> List[LangchainDocument]:
        """
        Process all supported files in a directory.

        Args:
            directory_path: Path to the directory

        Returns:
            List of LangchainDocument objects from all files
        """
        directory_path = Path(directory_path)

        if not directory_path.exists():
            raise FileNotFoundError(f"Directory not found: {directory_path}")

        all_documents = []

        for file_path in directory_path.rglob("*"):
            if (
                file_path.is_file()
                and file_path.suffix.lower() in self.supported_extensions
            ):
                try:
                    documents = self.process_file(str(file_path))
                    all_documents.extend(documents)
                except Exception as e:
                    logger.error(f"Error processing {file_path}: {str(e)}")
                    continue

        logger.info(f"Processed {len(all_documents)} total chunks from directory")
        return all_documents

    def _process_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files."""
        text_content = ""

        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text_content += f"\n--- Page {page_num + 1} ---\n{page_text}\n"

        return text_content

    def _process_docx(self, file_path: Path) -> str:
        """Extract text from DOCX files."""
        doc = Document(file_path)
        text_content = ""

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content += paragraph.text + "\n"

        return text_content

    def _process_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX files."""
        prs = Presentation(file_path)
        text_content = ""

        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"\n--- Slide {slide_num + 1} ---\n"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"

            if slide_text.strip():
                text_content += slide_text

        return text_content

    def _process_txt(self, file_path: Path) -> str:
        """Extract text from plain text files."""
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()

    def _process_image(self, file_path: Path) -> str:
        """Extract text from images using OCR."""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {str(e)}")
            return f"Image file: {file_path.name} (OCR processing failed)"

    def _process_csv(self, file_path: Path) -> str:
        """Extract text from CSV files."""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"CSV processing failed for {file_path}: {str(e)}")
            return f"CSV file: {file_path.name} (processing failed)"

    def clean_text(self, text: str) -> str:
        """
        Clean and normalize text content.

        Args:
            text: Raw text content

        Returns:
            Cleaned text content
        """
        # Remove extra whitespace
        text = re.sub(r"\s+", " ", text)

        # Remove special characters but keep punctuation
        text = re.sub(r"[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}]", "", text)

        # Normalize line breaks
        text = text.replace("\r\n", "\n").replace("\r", "\n")

        return text.strip()

    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats."""
        return list(self.supported_extensions.keys())
